import os
import pickle
import fitz  # PyMuPDF
from io import BytesIO
from striprtf.striprtf import rtf_to_text
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain_community.embeddings import OpenAIEmbeddings  # Updated import
from langchain_community.vectorstores import Chroma  # Updated import
from langchain.llms import OpenAI  # Ensure this is imported
from langchain.chains import RetrievalQA  # Ensure this is imported
import spacy
from nltk.corpus import stopwords
import nltk
from sklearn.decomposition import LatentDirichletAllocation
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import Counter
from keybert import KeyBERT
from textblob import TextBlob
from langdetect import detect
import textstat
from bertopic import BERTopic
from rake_nltk import Rake
from back_end.Utils.mongo_utils import MongoUtils
from transformers import pipeline
import time  # Ensure this is imported
import asyncio
import datetime
from datetime import timedelta
from slack_sdk import WebClient
from slack_sdk.errors import SlackApiError
import aiofiles
from langchain_community.embeddings import OpenAIEmbeddings  # Correct import
from sentence_transformers import SentenceTransformer

sentence_model = SentenceTransformer('all-MiniLM-L6-v2')

SLACK_BOT_TOKEN = os.getenv("OPENAI_API_KEY")

CHANNEL_IDS = ["C06HD8ADMC5"]  # Replace with actual channel IDs
client = WebClient(token=SLACK_BOT_TOKEN)

nltk.download('stopwords')

nlp = spacy.load("en_core_web_sm")
kw_model = KeyBERT()
topic_model = BERTopic()
ner_pipeline = pipeline("ner", model="dbmdz/bert-large-cased-finetuned-conll03-english")

rake = Rake()

from langdetect.lang_detect_exception import LangDetectException


class CustomTextLoader:
    """
    Custom loader for various document types.
    """

    def __init__(self):
        self.handlers = {
            '.docx': self.text_from_docx,
            '.pptx': self.text_from_pptx,
            '.pdf': self.text_from_pdf,
            '.rtf': self.text_from_rtf,
            '.txt': self.text_from_txt,
            '.csv': self.text_from_csv  # Added CSV handler
        }

    async def load(self, file, file_type=None, file_name=None):
        """
        Load text from a file based on its extension.
        """
        if file_type:
            ext = file_type
        elif file_name:
            ext = os.path.splitext(file_name)[-1].lower()
        else:
            ext = os.path.splitext(file.name if hasattr(file, 'name') else '')[-1].lower()

        print(f"Loading file with extension: {ext}")
        print(f"Handlers available: {list(self.handlers.keys())}")

        if ext in self.handlers:
            try:
                result = await self.handlers[ext](file, file_name)
                if ext == '.csv':
                    combined_text, initial_metadata, documents = result
                    return Document(page_content=combined_text, metadata=initial_metadata), documents
                else:
                    text, metadata = result
                    return Document(page_content=text, metadata=metadata)
            except Exception as e:
                print(f"Error processing {file_name or 'file'}: {e}")
                raise e
        else:
            raise ValueError(f"Unsupported file type {ext}")

    async def text_from_pdf(self, file, file_name=None):
        """
        Extract text from a PDF file.
        """
        print("Processing PDF file...")
        file.seek(0)
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ''
        for page in doc:
            text += page.get_text("text")
        doc.close()
        title = os.path.splitext(os.path.basename(file_name if file_name else ''))[0]
        return text, {"source": title}

    async def text_from_docx(self, file, file_name=None):
        """
        Extract text from a DOCX file.
        """
        print("Processing DOCX file...")
        import docx
        file.seek(0)
        doc = docx.Document(BytesIO(file.read()))
        text = '\n'.join(para.text for para in doc.paragraphs)
        title = os.path.splitext(os.path.basename(file_name if file_name else ''))[0]
        return text, {"source": title}

    async def text_from_pptx(self, file, file_name=None):
        """
        Extract text from a PPTX file.
        """
        print("Processing PPTX file...")
        import pptx
        file.seek(0)
        ppt = pptx.Presentation(BytesIO(file.read()))
        text = '\n'.join(
            shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
        )
        title = os.path.splitext(os.path.basename(file_name if file_name else ''))[0]
        return text, {"source": title}

    async def text_from_rtf(self, file, file_name=None):
        """
        Extract text from an RTF file.
        """
        print("Processing RTF file...")
        file.seek(0)
        text = rtf_to_text(file.read().decode("utf-8"))
        title = os.path.splitext(os.path.basename(file_name if file_name else ''))[0]
        return text, {"source": title}

    async def text_from_txt(self, file, file_name=None):
        """
        Extract text from a TXT file.
        """
        print("Processing TXT file...")
        file.seek(0)
        text = file.read().decode("utf-8")
        title = os.path.splitext(os.path.basename(file_name if file_name else ''))[0]
        return text, {"source": title}

    async def text_from_csv(self, file, file_name=None):
        """
        Extract text from a CSV file.
        """
        print("Processing CSV file...")
        file.seek(0)
        df = pd.read_csv(file)
        df = df.fillna('')
        company_name, sentiment = os.path.basename(file_name if file_name else '').rsplit('.', 1)[0].rsplit('_', 1)

        # Ensure the preprocess function is available here
        df['Processed_Text'] = df[sentiment].fillna("").astype(str).apply(preprocess)

        # Extract the reviews
        reviews = df['Processed_Text'].tolist()
        original_reviews = df[sentiment].tolist()

        # Use TF-IDF for vectorization
        vectorizer = TfidfVectorizer(max_features=5000)
        tfidf_matrix = vectorizer.fit_transform(reviews)

        # Use OpenAI embeddings
        openai_embeddings = OpenAIEmbeddings(model="text-embedding-ada-002")

        # Compute OpenAI embeddings
        openai_embed = openai_embeddings.embed_documents(reviews)
        openai_embed = np.array(openai_embed)  # Ensure embeddings are a NumPy array

        # Combine TF-IDF and OpenAI embeddings
        combined_embeddings = np.hstack((tfidf_matrix.toarray(), openai_embed))

        # Apply PCA for initial dimensionality reduction
        pca = PCA(n_components=100)
        pca_embeddings = pca.fit_transform(combined_embeddings)

        # Apply UMAP for further dimensionality reduction with increased components
        umap_model = UMAP(n_neighbors=50, n_components=20, min_dist=0.05, metric='cosine').fit(pca_embeddings)

        # Use HDBSCAN as the clustering algorithm with adjusted parameters
        hdbscan_model = hdbscan.HDBSCAN(min_cluster_size=5, min_samples=2, metric='euclidean',
                                        cluster_selection_epsilon=0.1, prediction_data=True)

        # Create a BERTopic model with custom embeddings, UMAP, and HDBSCAN clustering
        topic_model = BERTopic(embedding_model=openai_embeddings, umap_model=umap_model, hdbscan_model=hdbscan_model,
                               min_topic_size=2, calculate_probabilities=True)

        # Fit the model on the reviews
        topics, probabilities = topic_model.fit_transform(reviews, pca_embeddings)

        # Get topic info and representative documents
        topic_info = topic_model.get_topic_info()

        # Generate topic labels using the top 50 words for each topic
        topic_labels_dict = {}
        for topic in topic_info.Topic:
            if topic == -1:
                continue
            words = topic_model.get_topic(topic)
            if words:
                label = " ".join([word[0] for word in words[:50]])  # Get the top 50 words
                topic_labels_dict[topic] = label

        # Add a default label for outliers
        topic_labels_dict[-1] = "Outlier"

        # Ensure the topic_info DataFrame has the correct labels
        topic_info['Name'] = topic_info['Topic'].map(topic_labels_dict)

        # Map topic labels to the main DataFrame
        df['Topic'] = topics
        df['Topic_Name'] = df['Topic'].map(topic_labels_dict)

        # Create aggregated documents for each topic
        topic_groups = df.groupby('Topic')[sentiment].apply(lambda x: '_'.join(x)).reset_index()

        # Combine all texts for a single document
        combined_text = "\n".join(df['Processed_Text'].tolist())
        initial_metadata = {
            "source": os.path.basename(file_name if file_name else ''),
            "company name": company_name,
            "sentiment": sentiment,
            "topic": "",
            "topic_name": "",
            "categories": []
        }

        # Create Documents
        documents = []
        for i, row in topic_groups.iterrows():
            topic_number = row['Topic']
            doc = {
                "page_content": row[sentiment],
                "metadata": {
                    "source": "user_reviews",
                    "company name": company_name,
                    "sentiment": sentiment,
                    "topic": topic_number,
                    "topic_name": topic_labels_dict.get(topic_number, f"Topic {topic_number}")
                }
            }
            documents.append(Document(page_content=doc["page_content"], metadata=doc["metadata"]))

        return combined_text, initial_metadata, documents


def preprocess_text(text):
    """
    Preprocess and structure text.
    """
    # Removing excessive newlines and whitespace
    text = text.replace("\n\n", "\n").strip()
    # Split into lines
    lines = text.split("\n")

    structured_data = []
    slide = {}
    content = []
    for line in lines:
        if line.endswith('/12'):
            # New slide detected
            if slide:
                slide['content'] = "\n".join(content).strip()
                structured_data.append(slide)
                content = []
            slide = {"slide_number": line, "title": ""}
        elif any(keyword in line.lower() for keyword in
                 ["march", "avik bhandari", "switch to plum", "join our", "comprehensive", "easy claims"]):
            # Title or important section detected
            if content:
                slide['content'] = "\n".join(content).strip()
                structured_data.append(slide)
                content = []
            slide = {"slide_number": "", "title": line}
        else:
            content.append(line)

    # Append the last slide
    if slide:
        slide['content'] = "\n".join(content).strip()
        structured_data.append(slide)

    return structured_data


async def load_and_process_document(file_path, loader):
    """
    Load and split a document into chunks.
    """
    try:
        file_extension = os.path.splitext(file_path)[-1].lower()
        print(f"Loading file with extension: {file_extension}")

        async with aiofiles.open(file_path, 'rb') as file:
            file_content = await file.read()
            file_stream = BytesIO(file_content)

            # Load the document, handle the case when CSV is loaded
            if file_extension == '.csv':
                document_instance, documents = await loader.load(file_stream, file_type=file_extension,
                                                                 file_name=file_path)
                structured_data = documents
            else:
                document_instance = await loader.load(file_stream, file_type=file_extension, file_name=file_path)
                structured_data = preprocess_text(document_instance.page_content)
                documents = [Document(page_content=slide.get('content', ''),
                                      metadata={"source": os.path.basename(file_path),
                                                "slide_number": slide.get('slide_number', '')}) for slide in
                             structured_data]

            chunks = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100).split_documents(documents)
            print(f"Document {file_path} loaded and split into {len(chunks)} chunks.")
            for i, chunk in enumerate(chunks):
                print(f"Chunk {i + 1}: Length = {len(chunk.page_content)} characters")
            return chunks
    except Exception as e:
        print(f"Error processing {file_path}: {str(e)}")
        return []


async def embed_documents_in_batches(documents, embedding_model, batch_size=10, use_sentence_transformer=False):
    """
    Embed documents in batches for efficiency.
    """
    print("Embedding documents in batches...")
    batches = [documents[i:i + batch_size] for i in range(0, len(documents), batch_size)]
    embeddings = []

    loop = asyncio.get_event_loop()
    with ThreadPoolExecutor() as executor:
        if use_sentence_transformer:
            futures = [loop.run_in_executor(executor, embedding_model.encode, [doc.page_content for doc in batch], True)
                       for batch in batches]
        else:
            futures = [
                loop.run_in_executor(executor, embedding_model.embed_documents, [doc.page_content for doc in batch]) for
                batch in batches]

        results = await asyncio.gather(*futures)

        for result in results:
            if use_sentence_transformer:
                embeddings.extend(result.cpu().numpy())
            else:
                embeddings.extend(result)

    print(f"Completed embedding {len(documents)} documents.")
    return embeddings


def perform_bertopic_modeling(texts):
    """
    Perform topic modeling using BERTopic.
    """
    print("Performing BERTopic modeling...")
    if len(texts) > 1:
        topics, _ = topic_model.fit_transform(texts)
    else:
        topics = [0]  # Assign a default topic for single document cases
    print("BERTopic modeling completed.")
    return topics


def advanced_ner(content):
    """
    Perform advanced named entity recognition (NER).
    """
    ner_results = ner_pipeline(content)
    entities = {result['entity']: [] for result in ner_results}
    for result in ner_results:
        entities[result['entity']].append(result['word'])
    return entities


async def enhance_metadata_with_bertopic_and_advanced_ner(doc):
    content = doc.page_content
    topics = perform_bertopic_modeling([content])
    entities = advanced_ner(content)

    # Keyword Extraction
    keywords = kw_model.extract_keywords(content, keyphrase_ngram_range=(1, 2), stop_words='english', top_n=10)
    rake.extract_keywords_from_text(content)
    rake_keywords = rake.get_ranked_phrases_with_scores()[:10]

    # Sentiment Analysis
    text_blob = TextBlob(content)
    sentiment = text_blob.sentiment.polarity

    # Advanced Sentiment Analysis (using VADER)
    from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
    vader_analyzer = SentimentIntensityAnalyzer()
    vader_sentiment = vader_analyzer.polarity_scores(content)['compound']

    # Handle language detection with fallback
    try:
        language = detect(content)
    except LangDetectException as e:
        try:
            from fasttext import load_model
            lang_model = load_model('lid.176.bin')
            language = lang_model.predict(content)[0][0]
        except:
            language = "unknown"

    except LangDetectException:
        language = "unknown"

    # Readability Metrics
    flesch = textstat.flesch_kincaid_grade(content)
    smog = textstat.smog_index(content)
    readability_scores = {
        "flesch_kincaid_grade": flesch,
        "smog_index": smog,
        "flesch_reading_ease": textstat.flesch_reading_ease(content),
    }

    # Document Length
    length = len(content.split())

    ## Preserve original metadata
    original_metadata = {
        "source": doc.metadata.get("source", ""),
        "company name": doc.metadata.get("company name", ""),
        "sentiment": doc.metadata.get("sentiment", ""),
        "topic": doc.metadata.get("topic", ""),
        "topic_name": doc.metadata.get("topic_name", ""),
        "categories": doc.metadata.get("categories", [])
    }

    # Update metadata with new information, ensuring original metadata is preserved
    enhanced_metadata = {
        **doc.metadata,
        "entities": {key: ', '.join(map(str, val)) for key, val in entities.items()},
        "keywords": ', '.join([kw[0] for kw in keywords]),
        "rake_keywords": ', '.join([kw[1] for kw in rake_keywords]),
        "topics": ', '.join(map(str, topics)),  # Ensure topics are joined into a string
        "sentiment": str(sentiment),
        "vader_sentiment": str(vader_sentiment),
        "language": language,
        "readability_scores_flesch_kincaid_grade": str(readability_scores["flesch_kincaid_grade"]),
        "readability_scores_smog_index": str(readability_scores["smog_index"]),
        "readability_scores_flesch_reading_ease": str(readability_scores["flesch_reading_ease"]),
        "document_length": str(length),
    }
    # Combine original and enhanced metadata, prioritizing original values if present
    for key, value in original_metadata.items():
        if key not in enhanced_metadata or enhanced_metadata[key] == '':
            enhanced_metadata[key] = value

    doc.metadata = flatten_metadata(enhanced_metadata)
    print(doc.metadata)

    print(f"Metadata enhanced for document with source {doc.metadata['source']}")
    return doc


# Initialize MongoDB connection inside an async function
mongo_utils = MongoUtils()


async def initialize_mongo():
    db = await mongo_utils.connect_to_database()
    return db.embeddings


# Store Embeddings in MongoDB
async def store_embeddings(documents, embeddings):
    for doc, emb in zip(documents, embeddings):
        # Ensure embedding is a list
        embedding_list = emb.tolist() if not isinstance(emb, list) else emb
        doc.metadata['embedding'] = embedding_list  # Store embedding in metadata
        await embeddings_collection.insert_one({
            "content": doc.page_content,
            "metadata": doc.metadata,
        })


# Retrieve Embeddings from MongoDB
async def retrieve_embeddings():
    embeddings_data = await embeddings_collection.find({}).to_list(length=None)
    documents = [Document(page_content=data["content"], metadata=data["metadata"]) for data in embeddings_data]
    embeddings = [data["metadata"]["embedding"] for data in embeddings_data]
    return documents, embeddings


async def process_documents(main_document_path, file_paths, save_path):
    start_time = time.time()
    all_chunks = []
    loader = CustomTextLoader()
    openai_embedding_model = OpenAIEmbeddings(model="text-embedding-ada-002")

    global embeddings_collection
    embeddings_collection = await initialize_mongo()

    main_load_start = time.time()
    async with aiofiles.open(main_document_path, 'rb') as main_file:
        file_content = await main_file.read()
        result = await loader.load(BytesIO(file_content), file_type=os.path.splitext(main_document_path)[-1].lower(),
                                   file_name=main_document_path)
        if isinstance(result, tuple):
            main_doc_instance, _ = result
        else:
            main_doc_instance = result
        main_doc_embedding = openai_embedding_model.embed_documents([main_doc_instance.page_content])[0]
        print(f"Main document embedding length: {len(main_doc_embedding)}")
    print(f"Main document loaded and embedded in {time.time() - main_load_start} seconds")

    process_docs_start = time.time()
    tasks = [load_and_process_document(file_path, loader) for file_path in file_paths]
    results = await asyncio.gather(*tasks)

    for chunks in results:
        if chunks:
            all_chunks.extend(chunks)

    print(f"Other documents processed in {time.time() - process_docs_start} seconds")

    # Fetch and process Slack messages
    print("Fetching and processing Slack messages...")
    slack_documents, slack_re_embeddings = fetch_and_process_slack_messages()
    all_chunks.extend(slack_documents)
    print(f"Slack messages processed in {time.time() - process_docs_start} seconds")

    if all_chunks:
        tasks = [enhance_metadata_with_bertopic_and_advanced_ner(chunk) for chunk in all_chunks if
                 "slack_thread" not in chunk.metadata["source"]]
        all_chunks = await asyncio.gather(*tasks)

        similarity_start = time.time()
        chunk_embeddings = await embed_documents_in_batches(all_chunks, openai_embedding_model, batch_size=10)
        similarity_scores = [calculate_similarity(main_doc_embedding, chunk_embedding) for chunk_embedding in
                             chunk_embeddings]

        slack_similarity_scores = [calculate_similarity(main_doc_embedding, embedding) for embedding in
                                   slack_re_embeddings]

        all_similarities = similarity_scores + slack_similarity_scores
        all_documents = all_chunks + slack_documents
        doc_similarities = list(zip(all_documents, all_similarities))
        doc_similarities = sorted(doc_similarities, key=lambda x: x[1], reverse=True)
        print(f"Similarity calculations completed in {time.time() - similarity_start} seconds")

        await store_embeddings(all_documents, chunk_embeddings + slack_re_embeddings)
        print(f"Embeddings stored in MongoDB")

        print(f"Process documents completed in {time.time() - start_time} seconds")
        return None, doc_similarities, main_doc_embedding, all_chunks

    print(f"Process documents completed in {time.time() - start_time} seconds")
    return None, None, None, all_chunks


def fetch_messages(channel_id):
    messages = []
    try:
        result = client.conversations_history(channel=channel_id)
        for message in result['messages']:
            user_id = message.get('user', 'unknown')
            text = message.get('text', '')
            timestamp = datetime.fromtimestamp(float(message.get('ts', 0)))
            channel_info = client.conversations_info(channel=channel_id)
            channel_name = channel_info['channel']['name']
            user_info = client.users_info(user=user_id)
            user_name = user_info['user']['name']
            messages.append({
                "channel_id": channel_id,
                "channel_name": channel_name,
                "user_id": user_id,
                "user_name": user_name,
                "text": text,
                "timestamp": timestamp
            })
    except SlackApiError as e:
        print(f"Error fetching messages from channel {channel_id}: {e.response['error']}")
    return messages


def fetch_messages_from_slack():
    all_messages = []
    for channel_id in CHANNEL_IDS:
        all_messages.extend(fetch_messages(channel_id))
    return all_messages


def group_messages_into_threads(messages, time_threshold=timedelta(minutes=5)):
    """
    Group consecutive messages by the same user within a short time frame or by thread.
    """
    threads = []
    current_thread = []
    previous_message = None

    for message in messages:
        if previous_message and (
                message["user_id"] == previous_message["user_id"]
                and message["timestamp"] - previous_message["timestamp"] <= time_threshold
        ):
            current_thread.append(message)
        else:
            if current_thread:
                threads.append(current_thread)
            current_thread = [message]
        previous_message = message

    if current_thread:
        threads.append(current_thread)

    return threads


def flatten_metadata(metadata):
    """
    Convert complex metadata to a format that is compatible with MongoDB.
    """
    flattened_metadata = {}
    for key, value in metadata.items():
        if isinstance(value, dict):
            for sub_key, sub_value in value.items():
                flattened_metadata[f"{key}_{sub_key}"] = ', '.join(map(str, sub_value)) if isinstance(sub_value,
                                                                                                      list) else str(
                    sub_value)
        elif isinstance(value, list):
            flattened_metadata[key] = ', '.join(map(str, value))
        else:
            flattened_metadata[key] = str(value)
    return flattened_metadata


def process_threads_and_enhance_metadata(threads):
    """
    Process each thread to combine messages and enhance metadata.
    """
    enhanced_documents = []

    for thread in threads:
        combined_content = " ".join([msg["text"] for msg in thread])
        combined_metadata = {
            "source": "slack_thread",
            "user_ids": [msg["user_id"] for msg in thread],
            "user_names": [msg["user_name"] for msg in thread],
            "timestamps": [msg["timestamp"].strftime("%Y-%m-%d %H:%M:%S") for msg in thread],
            "channel_id": thread[0]["channel_id"],
            "channel_name": thread[0]["channel_name"],
        }

        doc = Document(page_content=combined_content, metadata=combined_metadata)
        enhanced_doc = enhance_metadata_with_bertopic_and_advanced_ner(doc)
        enhanced_documents.append(enhanced_doc)

    return enhanced_documents


def fetch_and_process_slack_messages():
    """
    Fetch Slack messages, group them into threads, and enhance metadata.
    """
    messages = fetch_messages_from_slack()
    threads = group_messages_into_threads(messages)
    enhanced_documents = process_threads_and_enhance_metadata(threads)

    # slack_embeddings = embed_documents_in_batches(enhanced_documents, sentence_model, batch_size=10, use_sentence_transformer=True)
    # slack_re_embeddings = embed_documents_in_batches(
    #     [Document(page_content=str(embedding)) for embedding in slack_embeddings],
    #     OpenAIEmbeddings(model="text-embedding-ada-002"),
    #     batch_size=10
    # )

    return enhanced_documents, enhanced_documents


def calculate_similarity(embedding1, embedding2, adjustment_factor=2):
    """
    Calculate similarity between two embeddings.
    """
    similarity = cosine_similarity([embedding1], [embedding2])[0][0]
    adjusted_similarity = similarity ** adjustment_factor
    return adjusted_similarity
