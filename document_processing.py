import os
import pickle
import fitz  # PyMuPDF
from io import BytesIO
from striprtf.striprtf import rtf_to_text
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain_community.embeddings import OpenAIEmbeddings  # Updated import
from langchain_community.vectorstores import Chroma  # Updated import
from langchain.llms import OpenAI  # Ensure this is imported
from langchain.chains import RetrievalQA  # Ensure this is imported
import spacy
from nltk.corpus import stopwords
import nltk
from sklearn.decomposition import LatentDirichletAllocation
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import Counter
from keybert import KeyBERT
from textblob import TextBlob
from langdetect import detect
import textstat
from bertopic import BERTopic
from rake_nltk import Rake
from back_end.Utils.mongo_utils import MongoUtils
from transformers import pipeline
import time  # Ensure this is imported

import datetime
from datetime import timedelta
from slack_sdk import WebClient
from slack_sdk.errors import SlackApiError
import os
from langchain_community.embeddings import OpenAIEmbeddings  # Correct import
from sentence_transformers import SentenceTransformer
sentence_model = SentenceTransformer('all-MiniLM-L6-v2')

SLACK_BOT_TOKEN = os.getenv("OPENAI_API_KEY")

CHANNEL_IDS = ["C06HD8ADMC5"]  # Replace with actual channel IDs
client = WebClient(token=SLACK_BOT_TOKEN)


nltk.download('stopwords')


nlp = spacy.load("en_core_web_sm")
kw_model = KeyBERT()
topic_model = BERTopic()
ner_pipeline = pipeline("ner", model="dbmdz/bert-large-cased-finetuned-conll03-english")

rake = Rake()

from langdetect.lang_detect_exception import LangDetectException

class CustomTextLoader:
    """
    Custom loader for various document types.
    """
    def __init__(self):
        self.handlers = {
            '.docx': self.text_from_docx,
            '.pptx': self.text_from_pptx,
            '.pdf': self.text_from_pdf,
            '.rtf': self.text_from_rtf,
            '.txt': self.text_from_txt,
        }

    def load(self, file):
        """
        Load text from a file based on its extension.
        """
        ext = os.path.splitext(file.name)[-1].lower()
        if ext in self.handlers:
            text = self.handlers[ext](file)
            title = os.path.splitext(os.path.basename(file.name))[0]
            return Document(page_content=text, metadata={"source": title})
        else:
            raise ValueError(f"Unsupported file type {ext}")

    def text_from_pdf(self, file):
        """
        Extract text from a PDF file.
        """
        file.seek(0)
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ''
        for page in doc:
            text += page.get_text("text")
        doc.close()
        return text

    def text_from_docx(self, file):
        """
        Extract text from a DOCX file.
        """
        import docx
        file.seek(0)
        doc = docx.Document(file)
        return '\n'.join(para.text for para in doc.paragraphs)

    def text_from_pptx(self, file):
        """
        Extract text from a PPTX file.
        """
        import pptx
        file.seek(0)
        ppt = pptx.Presentation(BytesIO(file.read()))
        return '\n'.join(
            shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
        )

    def text_from_rtf(self, file):
        """
        Extract text from an RTF file.
        """
        file.seek(0)
        return rtf_to_text(file.read().decode("utf-8"))

    def text_from_txt(self, file):
        """
        Extract text from a TXT file.
        """
        file.seek(0)
        return file.read().decode("utf-8")

def preprocess_text(text):
    """
    Preprocess and structure text.
    """
    # Removing excessive newlines and whitespace
    text = text.replace("\n\n", "\n").strip()
    # Split into lines
    lines = text.split("\n")
    
    structured_data = []
    slide = {}
    content = []
    for line in lines:
        if line.endswith('/12'):
            # New slide detected
            if slide:
                slide['content'] = "\n".join(content).strip()
                structured_data.append(slide)
                content = []
            slide = {"slide_number": line, "title": ""}
        elif any(keyword in line.lower() for keyword in ["march", "avik bhandari", "switch to plum", "join our", "comprehensive", "easy claims"]):
            # Title or important section detected
            if content:
                slide['content'] = "\n".join(content).strip()
                structured_data.append(slide)
                content = []
            slide = {"slide_number": "", "title": line}
        else:
            content.append(line)
    
    # Append the last slide
    if slide:
        slide['content'] = "\n".join(content).strip()
        structured_data.append(slide)
    
    return structured_data


def load_and_process_document(file_path, loader):
    """
    Load and split a document into chunks.
    """
    try:
        with open(file_path, 'rb') as file:
            document_instance = loader.load(file)
            # Preprocess and structure the text
            structured_text = preprocess_text(document_instance.page_content)
            # Save the parsed text to a file
            # Convert structured text back to Document objects
            documents = [Document(page_content=slide.get('content', ''), metadata={"source": os.path.basename(file_path), "slide_number": slide.get('slide_number', '')}) for slide in structured_text]
            chunks = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100).split_documents(documents)
            print(f"Document {file_path} loaded and split into {len(chunks)} chunks.")
            return chunks
    except Exception as e:
        print(f"Error processing {file_path}: {str(e)}")
        return []

def embed_documents_in_batches(documents, embedding_model, batch_size=10, use_sentence_transformer=False):
    """
    Embed documents in batches for efficiency.
    """
    print("Embedding documents in batches...")
    batches = [documents[i:i + batch_size] for i in range(0, len(documents), batch_size)]
    embeddings = []
    with ThreadPoolExecutor() as executor:
        if use_sentence_transformer:
            futures = [executor.submit(embedding_model.encode, [doc.page_content for doc in batch], convert_to_tensor=True) for batch in batches]
        else:
            futures = [executor.submit(embedding_model.embed_documents, [doc.page_content for doc in batch]) for batch in batches]
        for future in as_completed(futures):
            embeddings.extend(future.result().cpu().numpy() if use_sentence_transformer else future.result())
    print(f"Completed embedding {len(documents)} documents.")
    return embeddings

def perform_bertopic_modeling(texts):
    """
    Perform topic modeling using BERTopic.
    """
    print("Performing BERTopic modeling...")
    if len(texts) > 1:
        topics, _ = topic_model.fit_transform(texts)
    else:
        topics = [0]  # Assign a default topic for single document cases
    print("BERTopic modeling completed.")
    return topics

def advanced_ner(content):
    """
    Perform advanced named entity recognition (NER).
    """
    ner_results = ner_pipeline(content)
    entities = {result['entity']: [] for result in ner_results}
    for result in ner_results:
        entities[result['entity']].append(result['word'])
    return entities

def enhance_metadata_with_bertopic_and_advanced_ner(doc):
    """
    Enhance document metadata with BERTopic and advanced NER.
    """
    content = doc.page_content
    topics = perform_bertopic_modeling([content])
    entities = advanced_ner(content)
    
    # Keyword Extraction
    keywords = kw_model.extract_keywords(content, keyphrase_ngram_range=(1, 2), stop_words='english', top_n=10)
    rake.extract_keywords_from_text(content)
    rake_keywords = rake.get_ranked_phrases_with_scores()[:10]

    # Sentiment Analysis
    text_blob = TextBlob(content)
    sentiment = text_blob.sentiment.polarity
    
    # Advanced Sentiment Analysis (using VADER)
    from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
    vader_analyzer = SentimentIntensityAnalyzer()
    vader_sentiment = vader_analyzer.polarity_scores(content)['compound']
    
    # Handle language detection with fallback
    try:
        language = detect(content)
    except LangDetectException as e:
        try:
            from fasttext import load_model
            lang_model = load_model('lid.176.bin')
            language = lang_model.predict(content)[0][0]
        except:
            language = "unknown"
    
    # Readability Metrics
    flesch = textstat.flesch_kincaid_grade(content)
    smog = textstat.smog_index(content)
    readability_scores = {
        "flesch_kincaid_grade": flesch,
        "smog_index": smog,
        "flesch_reading_ease": textstat.flesch_reading_ease(content),
    }
    
    # Document Length
    length = len(content.split())
    
    # Update metadata
    doc.metadata.update({
        "entities": {key: ', '.join(map(str, val)) for key, val in entities.items()},
        "keywords": ', '.join([kw[0] for kw in keywords]),
        "rake_keywords": ', '.join([kw[1] for kw in rake_keywords]),
        "topics": topics,
        "sentiment": sentiment,
        "vader_sentiment": vader_sentiment,
        "language": language,
        "readability_scores": {k: str(v) for k, v in readability_scores.items()},
        "document_length": str(length),
    })

    # Flatten metadata
    doc.metadata = flatten_metadata(doc.metadata)

    print(f"Metadata enhanced for document with source {doc.metadata['source']}")
    return doc

# Initialize MongoDB
mongo_utils = MongoUtils()
db = mongo_utils.connect_to_database()
embeddings_collection = db.embeddings

# Store Embeddings in MongoDB
def store_embeddings(documents, embeddings):
    for doc, emb in zip(documents, embeddings):
        # Ensure embedding is a list
        embedding_list = emb.tolist() if not isinstance(emb, list) else emb
        doc.metadata['embedding'] = embedding_list  # Store embedding in metadata
        embeddings_collection.insert_one({
            "content": doc.page_content,
            "metadata": doc.metadata,
        })

        
# Retrieve Embeddings from MongoDB
def retrieve_embeddings():
    embeddings_data = list(embeddings_collection.find({}))
    documents = [Document(page_content=data["content"], metadata=data["metadata"]) for data in embeddings_data]
    embeddings = [data["metadata"]["embedding"] for data in embeddings_data]
    return documents, embeddings

def calculate_similarity(embedding1, embedding2, adjustment_factor=2):
    """
    Calculate similarity between two embeddings.
    """
    similarity = cosine_similarity([embedding1], [embedding2])[0][0]
    adjusted_similarity = similarity ** adjustment_factor
    return adjusted_similarity

def process_documents(main_document_path, file_paths, save_path):
    start_time = time.time()
    all_chunks = []
    loader = CustomTextLoader()
    openai_embedding_model = OpenAIEmbeddings(model="text-embedding-ada-002")

    main_load_start = time.time()
    with open(main_document_path, 'rb') as main_file:
        main_doc_instance = loader.load(main_file)
        main_doc_embedding = openai_embedding_model.embed_documents([main_doc_instance.page_content])[0]
    print(f"Main document loaded and embedded in {time.time() - main_load_start} seconds")

    process_docs_start = time.time()
    with ThreadPoolExecutor() as executor:
        futures = []
        for file_path in file_paths:
            futures.append(executor.submit(load_and_process_document, file_path, loader))

        for future in as_completed(futures):
            chunks = future.result()
            if chunks:
                all_chunks.extend(chunks)

    print(f"Other documents processed in {time.time() - process_docs_start} seconds")

    # Fetch and process Slack messages
    print("Fetching and processing Slack messages...")
    slack_documents, slack_re_embeddings = fetch_and_process_slack_messages()
    all_chunks.extend(slack_documents)
    print(f"Slack messages processed in {time.time() - process_docs_start} seconds")

    if all_chunks:
        all_chunks = [enhance_metadata_with_bertopic_and_advanced_ner(chunk) for chunk in all_chunks if "slack_thread" not in chunk.metadata["source"]]

        similarity_start = time.time()
        chunk_embeddings = embed_documents_in_batches(all_chunks, openai_embedding_model, batch_size=10)
        similarity_scores = [calculate_similarity(main_doc_embedding, chunk_embedding) for chunk_embedding in chunk_embeddings]
        
        slack_similarity_scores = [calculate_similarity(main_doc_embedding, embedding) for embedding in slack_re_embeddings]
        
        all_similarities = similarity_scores + slack_similarity_scores
        all_documents = all_chunks + slack_documents
        doc_similarities = list(zip(all_documents, all_similarities))
        doc_similarities = sorted(doc_similarities, key=lambda x: x[1], reverse=True)
        print(f"Similarity calculations completed in {time.time() - similarity_start} seconds")

        store_embeddings(all_documents, chunk_embeddings + slack_re_embeddings)
        print(f"Embeddings stored in MongoDB")

        print(f"Process documents completed in {time.time() - start_time} seconds")
        return None, doc_similarities, main_doc_embedding, all_chunks

    print(f"Process documents completed in {time.time() - start_time} seconds")
    return None, None, None, all_chunks


def fetch_messages(channel_id):
    messages = []
    try:
        result = client.conversations_history(channel=channel_id)
        for message in result['messages']:
            user_id = message.get('user', 'unknown')
            text = message.get('text', '')
            timestamp = datetime.fromtimestamp(float(message.get('ts', 0)))
            channel_info = client.conversations_info(channel=channel_id)
            channel_name = channel_info['channel']['name']
            user_info = client.users_info(user=user_id)
            user_name = user_info['user']['name']
            messages.append({
                "channel_id": channel_id,
                "channel_name": channel_name,
                "user_id": user_id,
                "user_name": user_name,
                "text": text,
                "timestamp": timestamp
            })
    except SlackApiError as e:
        print(f"Error fetching messages from channel {channel_id}: {e.response['error']}")
    return messages

def fetch_messages_from_slack():
    all_messages = []
    for channel_id in CHANNEL_IDS:
        all_messages.extend(fetch_messages(channel_id))
    return all_messages

def group_messages_into_threads(messages, time_threshold=timedelta(minutes=5)):
    """
    Group consecutive messages by the same user within a short time frame or by thread.
    """
    threads = []
    current_thread = []
    previous_message = None

    for message in messages:
        if previous_message and (
            message["user_id"] == previous_message["user_id"]
            and message["timestamp"] - previous_message["timestamp"] <= time_threshold
        ):
            current_thread.append(message)
        else:
            if current_thread:
                threads.append(current_thread)
            current_thread = [message]
        previous_message = message

    if current_thread:
        threads.append(current_thread)
    
    return threads

def flatten_metadata(metadata):
    """
    Convert complex metadata to a format that is compatible with MongoDB.
    """
    flattened_metadata = {}
    for key, value in metadata.items():
        if isinstance(value, dict):
            for sub_key, sub_value in value.items():
                flattened_metadata[f"{key}_{sub_key}"] = ', '.join(map(str, sub_value)) if isinstance(sub_value, list) else str(sub_value)
        elif isinstance(value, list):
            flattened_metadata[key] = ', '.join(map(str, value))
        else:
            flattened_metadata[key] = str(value)
    return flattened_metadata

def process_threads_and_enhance_metadata(threads):
    """
    Process each thread to combine messages and enhance metadata.
    """
    enhanced_documents = []

    for thread in threads:
        combined_content = " ".join([msg["text"] for msg in thread])
        combined_metadata = {
            "source": "slack_thread",
            "user_ids": [msg["user_id"] for msg in thread],
            "user_names": [msg["user_name"] for msg in thread],
            "timestamps": [msg["timestamp"].strftime("%Y-%m-%d %H:%M:%S") for msg in thread],
            "channel_id": thread[0]["channel_id"],
            "channel_name": thread[0]["channel_name"],
        }

        doc = Document(page_content=combined_content, metadata=combined_metadata)
        enhanced_doc = enhance_metadata_with_bertopic_and_advanced_ner(doc)
        enhanced_documents.append(enhanced_doc)

    return enhanced_documents

def fetch_and_process_slack_messages():
    """
    Fetch Slack messages, group them into threads, and enhance metadata.
    """
    messages = fetch_messages_from_slack()
    threads = group_messages_into_threads(messages)
    enhanced_documents = process_threads_and_enhance_metadata(threads)
    
    slack_embeddings = embed_documents_in_batches(enhanced_documents, sentence_model, batch_size=10, use_sentence_transformer=True)
    slack_re_embeddings = embed_documents_in_batches(
        [Document(page_content=str(embedding)) for embedding in slack_embeddings], 
        OpenAIEmbeddings(model="text-embedding-ada-002"), 
        batch_size=10
    )

    return enhanced_documents, slack_re_embeddings